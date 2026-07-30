"""
=============================================================================
  GuardianEye — PowerPoint Generator
  Generates a complete professional presentation for the FIM project
  Run: python make_ppt.py
  Output: GuardianEye_FIM_Presentation.pptx
=============================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette (dark cyber theme) ────────────────────────────────────────
C_BG_DARK    = RGBColor(0x0D, 0x11, 0x17)   # #0d1117 — slide background
C_BG_PANEL   = RGBColor(0x16, 0x1B, 0x22)   # #161b22 — card background
C_BG_CARD    = RGBColor(0x21, 0x26, 0x2D)   # #21262d — lighter card
C_ACCENT     = RGBColor(0x58, 0xA6, 0xFF)   # #58a6ff — blue accent
C_GREEN      = RGBColor(0x3F, 0xB9, 0x50)   # #3fb950 — green
C_YELLOW     = RGBColor(0xD2, 0x99, 0x22)   # #d29922 — yellow/warning
C_RED        = RGBColor(0xF8, 0x51, 0x49)   # #f85149 — red/danger
C_CYAN       = RGBColor(0x39, 0xC5, 0xCF)   # #39c5cf — cyan
C_WHITE      = RGBColor(0xE6, 0xED, 0xF3)   # #e6edf3 — primary text
C_GRAY       = RGBColor(0x8B, 0x94, 0x9E)   # #8b949e — secondary text
C_BORDER     = RGBColor(0x30, 0x36, 0x3D)   # #30363d — border

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# =============================================================================
# HELPERS
# =============================================================================

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Add a completely blank slide."""
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, colour=None):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour or C_BG_DARK


def add_rect(slide, x, y, w, h, fill_colour, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()   # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, colour=None,
                align=PP_ALIGN.LEFT, italic=False, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.name      = font_name
    run.font.color.rgb = colour or C_WHITE
    return txBox


def add_bullet_box(slide, lines, x, y, w, h,
                   font_size=16, title=None, title_colour=None, font_name="Segoe UI"):
    """Add a text box with optional title and multiple bullet lines."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in lines:
        if first and title is None:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        colour = C_WHITE
        sz     = font_size
        bd     = False
        prefix = ""

        if isinstance(line, tuple):
            text, colour, sz, bd = (list(line) + [font_size, False])[:4]
        else:
            text = line

        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(sz)
        run.font.bold      = bd
        run.font.name      = font_name
        run.font.color.rgb = colour if isinstance(colour, RGBColor) else C_WHITE
        p.space_after      = Pt(4)
        first = False

    return txBox


def add_code_box(slide, code_text, x, y, w, h, font_size=13):
    """Dark monospaced code box."""
    rect = add_rect(slide, x, y, w, h, C_BG_CARD)
    rect.line.color.rgb = C_BORDER
    rect.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.1),
        w - Inches(0.3), h - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = False

    first = True
    for line in code_text.split("\n"):
        if first:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text           = line
        run.font.size      = Pt(font_size)
        run.font.name      = "Consolas"
        run.font.color.rgb = C_CYAN
        first = False

    return txBox


def slide_header(slide, title, subtitle=None, accent_bar=True):
    """Add consistent title + optional subtitle + accent bar at top."""
    if accent_bar:
        add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), C_ACCENT)

    add_textbox(slide, title,
                Inches(0.6), Inches(0.15), Inches(11), Inches(0.7),
                font_size=28, bold=True, colour=C_WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.6), Inches(0.8), Inches(11), Inches(0.4),
                    font_size=14, colour=C_GRAY)


def card(slide, x, y, w, h, title, title_colour=None):
    """Draw a card with a coloured header bar."""
    add_rect(slide, x, y, w, h, C_BG_PANEL)
    add_rect(slide, x, y, w, Inches(0.04), title_colour or C_ACCENT)
    add_textbox(slide, title,
                x + Inches(0.15), y + Inches(0.08),
                w - Inches(0.3), Inches(0.4),
                font_size=13, bold=True, colour=title_colour or C_ACCENT)
    return y + Inches(0.52)   # return content start Y


# =============================================================================
# SLIDES
# =============================================================================

def slide_01_title(prs):
    """Title slide."""
    s = blank_slide(prs)
    fill_bg(s)

    # Full-width accent bar bottom
    add_rect(s, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), C_ACCENT)

    # Left accent strip
    add_rect(s, Inches(0), Inches(0), Inches(0.08), SLIDE_H, C_ACCENT)

    # Shield circle
    circle = s.shapes.add_shape(9, Inches(1.0), Inches(1.2),
                                 Inches(1.8), Inches(1.8))
    circle.fill.solid();  circle.fill.fore_color.rgb = C_ACCENT
    circle.line.fill.background()

    add_textbox(s, "G",
                Inches(1.55), Inches(1.55), Inches(0.7), Inches(0.8),
                font_size=36, bold=True, colour=C_BG_DARK,
                align=PP_ALIGN.CENTER, font_name="Segoe UI")

    # Main title
    add_textbox(s, "GuardianEye",
                Inches(1.0), Inches(3.1), Inches(10), Inches(1.1),
                font_size=54, bold=True, colour=C_WHITE, font_name="Segoe UI")

    add_textbox(s, "File Integrity Monitor",
                Inches(1.0), Inches(4.15), Inches(10), Inches(0.7),
                font_size=30, colour=C_ACCENT, font_name="Segoe UI")

    add_textbox(s, "Real-time file monitoring  •  SHA-256 hashing  •  Instant threat detection",
                Inches(1.0), Inches(4.9), Inches(10), Inches(0.5),
                font_size=15, colour=C_GRAY, font_name="Segoe UI")

    add_textbox(s, "Built with Python  |  Standard Library Only  |  GUI + CLI",
                Inches(0.6), Inches(7.05), Inches(10), Inches(0.35),
                font_size=12, colour=C_BG_DARK, bold=True,
                font_name="Segoe UI")


def slide_02_problem(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "The Problem", "How do you know if your files have been tampered with?")

    # Three problem cards
    problems = [
        ("🔴  Silent Modification",
         C_RED,
         "An attacker edits a config file or policy document.\nNo notification. No log. No one knows."),
        ("🔴  Malware Drops",
         C_YELLOW,
         "A malicious file is quietly placed inside a\nprotected directory. Completely undetected."),
        ("🔴  Insider Threats",
         C_RED,
         "A rogue employee deletes critical files or\ninjects fraudulent records into a database CSV."),
    ]

    for i, (title, col, desc) in enumerate(problems):
        x = Inches(0.5 + i * 4.2)
        content_y = card(s, x, Inches(1.5), Inches(3.9), Inches(4.8), title, col)
        add_textbox(s, desc,
                    x + Inches(0.15), content_y, Inches(3.6), Inches(3.5),
                    font_size=14, colour=C_GRAY)

    add_textbox(s,
        "Traditional monitoring tools check logs — but logs can be deleted too.\n"
        "The only reliable method is to fingerprint the files themselves.",
        Inches(0.5), Inches(6.5), Inches(12), Inches(0.7),
        font_size=13, colour=C_ACCENT, align=PP_ALIGN.CENTER)


def slide_03_what_is_fim(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "What is File Integrity Monitoring?",
                 "FIM detects unauthorised changes to files by comparing cryptographic hashes")

    add_textbox(s, "Core Concept",
                Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.4),
                font_size=16, bold=True, colour=C_ACCENT)

    steps = [
        ("STEP 1",  C_GREEN,  "Fingerprint every file with SHA-256"),
        ("STEP 2",  C_ACCENT, "Store all  filepath → hash  pairs as a baseline"),
        ("STEP 3",  C_ACCENT, "Re-fingerprint files continuously"),
        ("STEP 4",  C_YELLOW, "Compare new hashes with the baseline"),
        ("STEP 5",  C_RED,    "Alert immediately if any mismatch is found"),
    ]

    for i, (label, col, text) in enumerate(steps):
        y = Inches(2.0 + i * 0.88)
        add_rect(s, Inches(0.5), y, Inches(1.2), Inches(0.65), col)
        add_textbox(s, label,
                    Inches(0.5), y + Inches(0.1), Inches(1.2), Inches(0.5),
                    font_size=11, bold=True, colour=C_BG_DARK, align=PP_ALIGN.CENTER)
        add_textbox(s, text,
                    Inches(1.9), y + Inches(0.12), Inches(4.5), Inches(0.5),
                    font_size=15, colour=C_WHITE)

    # Hash example
    content_y = card(s, Inches(7.2), Inches(1.5), Inches(5.7), Inches(5.5),
                     "SHA-256 Example", C_CYAN)
    add_code_box(s,
        "Original file:\n"
        "\"admin  ALLOW ALL\"\n\n"
        "SHA-256:\nf715b9eeb9260d91f6...\n\n"
        "─────────────────────\n\n"
        "After 1 character changed:\n"
        "\"admin  ALLOW ALL \"\n"
        "         (added space)\n\n"
        "SHA-256:\n30a827acf5634be6b3...\n\n"
        "→ Completely different hash\n"
        "→ MODIFIED alert fires",
        Inches(7.35), content_y, Inches(5.4), Inches(4.7), font_size=11)


def slide_04_solution(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Our Solution — GuardianEye",
                 "A Python-based FIM tool with GUI + CLI, built entirely on the standard library")

    features = [
        ("✔  Real-time monitoring",     C_GREEN,  "Checks files every 1 second continuously"),
        ("✔  SHA-256 fingerprinting",   C_GREEN,  "Cryptographically strong — 1 byte change = new hash"),
        ("✔  Three detection types",    C_CYAN,   "Modified / Deleted / New File"),
        ("✔  Live GUI alerts",          C_ACCENT, "Colour-coded feed with timestamps"),
        ("✔  Audit log",               C_ACCENT, "guardian_events.log — every event timestamped"),
        ("✔  Trust store (JSON)",       C_CYAN,   "Human-readable — inspect hashes in Notepad"),
        ("✔  Zero dependencies",        C_GREEN,  "Standard Python only — no pip install needed"),
        ("✔  CLI + GUI modes",          C_YELLOW, "fim_scanner.py (terminal) + guardian_gui.py (GUI)"),
    ]

    for i, (feat, col, desc) in enumerate(features):
        col_n = i % 2
        row_n = i // 2
        x = Inches(0.5 + col_n * 6.4)
        y = Inches(1.5 + row_n * 1.3)
        add_rect(s, x, y, Inches(6.0), Inches(1.1), C_BG_PANEL)
        add_rect(s, x, y, Inches(0.05), Inches(1.1), col)
        add_textbox(s, feat,
                    x + Inches(0.2), y + Inches(0.05), Inches(5.7), Inches(0.45),
                    font_size=14, bold=True, colour=col)
        add_textbox(s, desc,
                    x + Inches(0.2), y + Inches(0.5), Inches(5.7), Inches(0.45),
                    font_size=12, colour=C_GRAY)


def slide_05_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "System Architecture",
                 "Two-layer design — core engine + GUI wrapper")

    # GUI layer box
    add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.6), C_BG_PANEL)
    add_textbox(s, "GUI Layer  —  guardian_gui.py  (tkinter)",
                Inches(0.7), Inches(1.45), Inches(6), Inches(0.45),
                font_size=13, bold=True, colour=C_ACCENT)

    gui_boxes = [
        ("Sidebar\nControls", C_ACCENT),
        ("Live Alert\nFeed", C_YELLOW),
        ("Enrolled\nFiles Tab", C_GREEN),
        ("Event\nLog Tab", C_CYAN),
        ("WatcherThread\n(background)", C_RED),
    ]
    for i, (lbl, col) in enumerate(gui_boxes):
        x = Inches(0.65 + i * 2.4)
        add_rect(s, x, Inches(1.95), Inches(2.1), Inches(0.9), col)
        add_textbox(s, lbl, x + Inches(0.05), Inches(2.0), Inches(2.0), Inches(0.85),
                    font_size=11, bold=True, colour=C_BG_DARK, align=PP_ALIGN.CENTER)

    # Arrow down
    add_textbox(s, "⬇  imports functions from",
                Inches(5.5), Inches(4.1), Inches(3), Inches(0.4),
                font_size=12, colour=C_GRAY, align=PP_ALIGN.CENTER)

    # Core engine box
    add_rect(s, Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.5), C_BG_CARD)
    add_textbox(s, "Core Engine  —  fim_scanner.py",
                Inches(0.7), Inches(4.6), Inches(6), Inches(0.45),
                font_size=13, bold=True, colour=C_GREEN)

    core_funcs = [
        ("fingerprint()", "SHA-256 hash a file", C_GREEN),
        ("snapshot_directory()", "Hash all files in folder", C_CYAN),
        ("save_trust_store()", "Write trust_store.json", C_ACCENT),
        ("load_trust_store()", "Read trust_store.json", C_ACCENT),
        ("record_event()", "Append to events.log", C_YELLOW),
    ]
    for i, (fn, desc, col) in enumerate(core_funcs):
        x = Inches(0.7 + i * 2.45)
        add_rect(s, x, Inches(5.15), Inches(2.25), Inches(0.8), C_BG_PANEL)
        add_textbox(s, fn, x + Inches(0.08), Inches(5.18), Inches(2.1), Inches(0.35),
                    font_size=10, bold=True, colour=col, font_name="Consolas")
        add_textbox(s, desc, x + Inches(0.08), Inches(5.5), Inches(2.1), Inches(0.35),
                    font_size=9, colour=C_GRAY)

    # Files row
    for i, (fname, col) in enumerate([
        ("watched_files/", C_CYAN),
        ("trust_store.json", C_GREEN),
        ("guardian_events.log", C_YELLOW),
    ]):
        x = Inches(1.5 + i * 3.8)
        add_textbox(s, fname,
                    x, Inches(7.1), Inches(3.5), Inches(0.35),
                    font_size=12, colour=col, font_name="Consolas",
                    align=PP_ALIGN.CENTER)


def slide_06_how_it_works(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "How It Works — Detection Loop",
                 "WatcherThread runs every 1 second in the background")

    steps = [
        ("1", "snapshot_directory()",
         "Walk watched_files/ recursively\nHash every file with SHA-256\nReturn dict {filepath: {sha256, size}}",
         C_CYAN),
        ("2", "Compare against baseline",
         "For each live file:\n  • In baseline + hash changed → MODIFIED\n  • Not in baseline at all → NEW FILE\nFor each baseline file:\n  • Not found on disk → DELETED",
         C_YELLOW),
        ("3", "Post to queue",
         "alert_queue.put((\"alert\", kind, path, detail))\nrecord_event() → guardian_events.log\nGUI polls queue every 200ms safely",
         C_RED),
        ("4", "Update baseline",
         "self.trusted = dict(live)\nCurrent state becomes new reference\nSame alert NEVER fires twice",
         C_GREEN),
    ]

    for i, (num, title, body, col) in enumerate(steps):
        x = Inches(0.4 + i * 3.2)
        # Number circle
        circle = s.shapes.add_shape(9, x + Inches(1.2), Inches(1.4),
                                     Inches(0.7), Inches(0.7))
        circle.fill.solid(); circle.fill.fore_color.rgb = col
        circle.line.fill.background()
        add_textbox(s, num, x + Inches(1.2), Inches(1.45),
                    Inches(0.7), Inches(0.55),
                    font_size=18, bold=True, colour=C_BG_DARK,
                    align=PP_ALIGN.CENTER)

        # Arrow between steps
        if i < 3:
            add_textbox(s, "→",
                        x + Inches(3.0), Inches(1.5), Inches(0.4), Inches(0.5),
                        font_size=20, colour=C_GRAY, align=PP_ALIGN.CENTER)

        add_rect(s, x, Inches(2.2), Inches(3.0), Inches(4.8), C_BG_PANEL)
        add_rect(s, x, Inches(2.2), Inches(3.0), Inches(0.04), col)
        add_textbox(s, title,
                    x + Inches(0.15), Inches(2.28), Inches(2.7), Inches(0.45),
                    font_size=13, bold=True, colour=col, font_name="Consolas")
        add_textbox(s, body,
                    x + Inches(0.15), Inches(2.78), Inches(2.7), Inches(4.0),
                    font_size=12, colour=C_GRAY)


def slide_07_detection_types(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Three Detection Types",
                 "Every type of file system change is caught and classified")

    types = [
        ("MODIFIED", C_YELLOW, "~",
         "File content has changed since the baseline was taken.\n\n"
         "Example:\n"
         "  • security_policy.txt weakened\n"
         "  • firewall_rules.txt changed to ALLOW ALL\n"
         "  • employee_data.csv row injected\n\n"
         "Detection:\nSHA-256 of live file ≠ SHA-256 in baseline\n\n"
         "Shows:\nOld hash vs new hash side by side"),
        ("DELETED", C_RED, "−",
         "A file from the baseline no longer exists on disk.\n\n"
         "Example:\n"
         "  • api_keys.txt deleted after exfiltration\n"
         "  • checksums.txt removed to hide tampering\n"
         "  • audit_log.csv deleted by insider\n\n"
         "Detection:\nBaseline has the path but os.path.exists() = False\n\n"
         "Shows:\nFilename + 'file no longer exists on disk'"),
        ("NEW FILE", C_CYAN, "+",
         "A file appeared that was not in the baseline.\n\n"
         "Example:\n"
         "  • malware_payload.exe.txt dropped\n"
         "  • ransomware_note.txt created\n"
         "  • Backdoor script added to folder\n\n"
         "Detection:\nLive file path not found in baseline dict\n\n"
         "Shows:\nFilename + SHA-256 hash of new file"),
    ]

    for i, (label, col, icon, body) in enumerate(types):
        x = Inches(0.4 + i * 4.3)
        # Icon circle
        circle = s.shapes.add_shape(9, x + Inches(1.3), Inches(1.4),
                                     Inches(0.8), Inches(0.8))
        circle.fill.solid(); circle.fill.fore_color.rgb = col
        circle.line.fill.background()
        add_textbox(s, icon, x + Inches(1.3), Inches(1.45),
                    Inches(0.8), Inches(0.65),
                    font_size=26, bold=True, colour=C_BG_DARK,
                    align=PP_ALIGN.CENTER)

        add_rect(s, x, Inches(2.35), Inches(4.0), Inches(4.8), C_BG_PANEL)
        add_rect(s, x, Inches(2.35), Inches(4.0), Inches(0.05), col)
        add_textbox(s, f"[{icon}]  {label}",
                    x + Inches(0.15), Inches(2.43), Inches(3.7), Inches(0.5),
                    font_size=16, bold=True, colour=col)
        add_textbox(s, body,
                    x + Inches(0.15), Inches(2.98), Inches(3.7), Inches(3.8),
                    font_size=12, colour=C_GRAY)


def slide_08_gui(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "GUI Overview — guardian_gui.py",
                 "Dark SOC-style interface built entirely with Python tkinter")

    sections = [
        ("Header Bar", C_ACCENT,
         "• GuardianEye logo + title\n• Status dot: grey=Idle, green=Watching"),
        ("Left Sidebar", C_GREEN,
         "• Trust store status card\n• 🔒 Build Trust Store button\n"
         "• ▶ Start Watching button\n• ■ Stop Watching button\n"
         "• Monitored folder path\n• Session stats (Modified/Deleted/New)\n"
         "• 📂 Open Folder  📄 View Log"),
        ("Live Alerts Tab", C_YELLOW,
         "• Real-time colour-coded feed\n• Yellow = Modified\n"
         "• Red = Deleted\n• Cyan = New File\n"
         "• Old hash vs new hash shown\n• Alert counter in toolbar"),
        ("Enrolled Files Tab", C_CYAN,
         "• Table of all fingerprinted files\n• Columns: File / SHA-256 / Size / Time\n"
         "• Refresh button to reload"),
        ("Event Log Tab", C_RED,
         "• Full guardian_events.log contents\n• Colour highlighted same as alerts\n"
         "• Horizontal scroll for long hashes"),
        ("Status Bar", C_GRAY,
         "• Live status message\n• Real-time clock (updates every second)"),
    ]

    for i, (title, col, body) in enumerate(sections):
        col_n = i % 3
        row_n = i // 3
        x = Inches(0.4 + col_n * 4.3)
        y = Inches(1.5 + row_n * 2.7)
        add_rect(s, x, y, Inches(4.0), Inches(2.5), C_BG_PANEL)
        add_rect(s, x, y, Inches(4.0), Inches(0.05), col)
        add_textbox(s, title,
                    x + Inches(0.15), y + Inches(0.1), Inches(3.7), Inches(0.4),
                    font_size=13, bold=True, colour=col)
        add_textbox(s, body,
                    x + Inches(0.15), y + Inches(0.55), Inches(3.7), Inches(1.8),
                    font_size=12, colour=C_GRAY)


def slide_09_tech_stack(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Technology Stack",
                 "100% Python standard library — no external packages required")

    modules = [
        ("hashlib",    C_GREEN,  "SHA-256 / MD5\nCryptographic fingerprinting of files"),
        ("tkinter",    C_ACCENT, "GUI Framework\nFull GUI — windows, buttons, tabs, text boxes"),
        ("threading",  C_YELLOW, "Background Thread\nWatcher runs without blocking the UI"),
        ("queue",      C_CYAN,   "Thread-safe Queue\nPasses alerts from watcher → GUI safely"),
        ("json",       C_GREEN,  "Trust Store\nHuman-readable baseline storage format"),
        ("os / shutil",C_ACCENT, "File System\nWalk directories, check paths, restore files"),
        ("time",       C_GRAY,   "Polling Interval\n1-second sleep between each scan cycle"),
        ("datetime",   C_YELLOW, "Timestamps\nAll alerts and log entries are time-stamped"),
    ]

    for i, (name, col, desc) in enumerate(modules):
        col_n = i % 4
        row_n = i // 4
        x = Inches(0.4 + col_n * 3.2)
        y = Inches(1.5 + row_n * 2.7)
        add_rect(s, x, y, Inches(2.95), Inches(2.5), C_BG_PANEL)
        add_rect(s, x, y, Inches(2.95), Inches(0.05), col)
        add_textbox(s, name,
                    x + Inches(0.12), y + Inches(0.1), Inches(2.7), Inches(0.45),
                    font_size=16, bold=True, colour=col, font_name="Consolas")
        add_textbox(s, desc,
                    x + Inches(0.12), y + Inches(0.6), Inches(2.7), Inches(1.7),
                    font_size=12, colour=C_GRAY)


def slide_10_file_structure(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Project File Structure",
                 "Clean two-file architecture — engine + GUI")

    add_code_box(s,
        "g:\\fim_tool\\\n"
        "│\n"
        "├── fim_scanner.py          ← Core engine\n"
        "│     fingerprint()         → SHA-256 hash a file\n"
        "│     snapshot_directory()  → hash all files\n"
        "│     save_trust_store()    → write JSON baseline\n"
        "│     load_trust_store()    → read JSON baseline\n"
        "│     record_event()        → append to log\n"
        "│\n"
        "├── guardian_gui.py         ← GUI layer (tkinter)\n"
        "│     GuardianEyeApp        → main window class\n"
        "│     WatcherThread         → background monitor\n"
        "│\n"
        "├── Start-SentinelFIM.bat   ← Double-click launcher\n"
        "│\n"
        "├── watched_files\\          ← Monitored folder\n"
        "│   ├── a.txt\n"
        "│   ├── b.txt  c.txt  d.txt\n"
        "│   └── sensitive\\credentials.txt\n"
        "│\n"
        "├── trust_store.json        ← SHA-256 baseline\n"
        "└── guardian_events.log     ← Audit trail",
        Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.8), font_size=12)

    # trust_store.json example
    content_y = card(s, Inches(7.0), Inches(1.4), Inches(5.9), Inches(5.8),
                     "trust_store.json format", C_GREEN)
    add_code_box(s,
        "{\n"
        "  \"guardian_eye\": \"1.0\",\n"
        "  \"captured_at\": \"2026-07-30T16:00:00\",\n"
        "  \"total_files\": 5,\n"
        "  \"fingerprints\": {\n"
        "    \"a.txt\": {\n"
        "      \"sha256\": \"f715b9eeb926...\",\n"
        "      \"size\": 108,\n"
        "      \"seen\": \"2026-07-30 16:00:00\"\n"
        "    },\n"
        "    \"b.txt\": {\n"
        "      \"sha256\": \"22bed23318...\",\n"
        "      \"size\": 81,\n"
        "      \"seen\": \"2026-07-30 16:00:00\"\n"
        "    }\n"
        "  }\n"
        "}",
        Inches(7.15), content_y, Inches(5.6), Inches(4.7), font_size=12)


def slide_11_how_to_test(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "How to Test — Live Demo Steps",
                 "Two terminals: one watching, one tampering")

    steps = [
        ("STEP 1", C_GREEN,  "Build Trust Store",
         "Click 🔒 Build Trust Store in the GUI\n"
         "→ Restores 5 clean files\n"
         "→ SHA-256 fingerprints all of them\n"
         "→ Saves trust_store.json"),
        ("STEP 2", C_ACCENT, "Start Watching",
         "Click ▶ Start Watching\n"
         "→ Takes live snapshot as baseline\n"
         "→ Zero alerts at startup\n"
         "→ Status bar: 'Watching 5 file(s)'"),
        ("STEP 3", C_YELLOW, "Tamper Files",
         "Open watched_files/ in Explorer\n"
         "→ Edit a.txt → MODIFIED alert\n"
         "→ Delete b.txt → DELETED alert\n"
         "→ Create evil.txt → NEW FILE alert"),
        ("STEP 4", C_RED,    "See Live Alerts",
         "Within 1 second of each change:\n"
         "→ Colour-coded alert appears\n"
         "→ Old vs new hash shown\n"
         "→ Event logged to guardian_events.log"),
    ]

    for i, (label, col, title, body) in enumerate(steps):
        x = Inches(0.4 + i * 3.2)
        # Step number
        add_rect(s, x, Inches(1.4), Inches(2.9), Inches(0.55), col)
        add_textbox(s, label,
                    x, Inches(1.43), Inches(2.9), Inches(0.45),
                    font_size=13, bold=True, colour=C_BG_DARK,
                    align=PP_ALIGN.CENTER)
        add_rect(s, x, Inches(1.95), Inches(2.9), Inches(4.9), C_BG_PANEL)
        add_textbox(s, title,
                    x + Inches(0.15), Inches(2.02), Inches(2.6), Inches(0.5),
                    font_size=14, bold=True, colour=col)
        add_textbox(s, body,
                    x + Inches(0.15), Inches(2.6), Inches(2.6), Inches(3.8),
                    font_size=13, colour=C_GRAY)

    add_textbox(s,
        "Tip:  python tamper_demo.py  — automatically simulates 5 attack types with countdown",
        Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
        font_size=12, colour=C_CYAN, align=PP_ALIGN.CENTER, font_name="Consolas")


def slide_12_use_cases(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Use Cases & Applications",
                 "Where GuardianEye fits in real-world security scenarios")

    cases = [
        ("🏢  Enterprise Security",   C_ACCENT,
         "Monitor critical config files,\nfirewall rules, and access\ncontrol lists for unauthorised changes."),
        ("🏥  Healthcare / HIPAA",    C_GREEN,
         "Ensure patient records and\ncompliance documents are not\nmodified or deleted without authorisation."),
        ("🖥️  Server Hardening",       C_CYAN,
         "Watch system files, SSH keys,\nand web server configs for\nany modification after deployment."),
        ("🔬  SOC / Blue Team Labs",  C_YELLOW,
         "Demonstrate FIM concepts,\npractice incident detection,\nand build forensic audit trails."),
        ("🎓  Cybersecurity Learning", C_RED,
         "Understand how SHA-256 works,\nhow FIM tools detect threats,\nand how GUI tools are built in Python."),
        ("🔐  Insider Threat Detection", C_ACCENT,
         "Catch data exfiltration attempts\nwhere files are modified or\ndeleted by internal users."),
    ]

    for i, (title, col, body) in enumerate(cases):
        col_n = i % 3
        row_n = i // 3
        x = Inches(0.4 + col_n * 4.3)
        y = Inches(1.5 + row_n * 2.7)
        add_rect(s, x, y, Inches(4.0), Inches(2.5), C_BG_PANEL)
        add_rect(s, x, y, Inches(4.0), Inches(0.05), col)
        add_textbox(s, title,
                    x + Inches(0.15), y + Inches(0.1), Inches(3.7), Inches(0.5),
                    font_size=13, bold=True, colour=col)
        add_textbox(s, body,
                    x + Inches(0.15), y + Inches(0.65), Inches(3.7), Inches(1.7),
                    font_size=12, colour=C_GRAY)


def slide_13_conclusion(prs):
    s = blank_slide(prs)
    fill_bg(s)

    add_rect(s, Inches(0), Inches(0), Inches(0.08), SLIDE_H, C_ACCENT)
    add_rect(s, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), C_BG_PANEL)

    add_textbox(s, "GuardianEye",
                Inches(1.0), Inches(1.2), Inches(11), Inches(1.1),
                font_size=48, bold=True, colour=C_WHITE)
    add_textbox(s, "File Integrity Monitor — Project Summary",
                Inches(1.0), Inches(2.25), Inches(11), Inches(0.6),
                font_size=22, colour=C_ACCENT)

    summary = [
        "✔   Built entirely in Python standard library — no pip install",
        "✔   SHA-256 cryptographic fingerprinting for tamper detection",
        "✔   Detects Modified, Deleted, and New files in real time",
        "✔   Dark SOC-style GUI with live colour-coded alert feed",
        "✔   Permanent audit log with timestamps for forensics",
        "✔   Alerts fire exactly once per change — no repeated noise",
        "✔   CLI mode (fim_scanner.py) + GUI mode (guardian_gui.py)",
    ]

    for i, line in enumerate(summary):
        add_textbox(s, line,
                    Inches(1.0), Inches(3.05 + i * 0.5), Inches(11), Inches(0.45),
                    font_size=14, colour=C_WHITE if i % 2 == 0 else C_GRAY)

    add_textbox(s,
        "github.com/joshmadakor1/PowerShell-Integrity-FIM  (original concept)  "
        "→  Ported & extended to Python with GUI",
        Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
        font_size=11, colour=C_GRAY, align=PP_ALIGN.CENTER)


# =============================================================================
# MAIN — build the full deck
# =============================================================================

def build_presentation():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs)         ; print("  [1/13] Title")
    slide_02_problem(prs)       ; print("  [2/13] Problem Statement")
    slide_03_what_is_fim(prs)   ; print("  [3/13] What is FIM?")
    slide_04_solution(prs)      ; print("  [4/13] Our Solution")
    slide_05_architecture(prs)  ; print("  [5/13] Architecture")
    slide_06_how_it_works(prs)  ; print("  [6/13] Detection Loop")
    slide_07_detection_types(prs); print("  [7/13] Detection Types")
    slide_08_gui(prs)           ; print("  [8/13] GUI Overview")
    slide_09_tech_stack(prs)    ; print("  [9/13] Tech Stack")
    slide_10_file_structure(prs); print("  [10/13] File Structure")
    slide_11_how_to_test(prs)   ; print("  [11/13] How to Test")
    slide_12_use_cases(prs)     ; print("  [12/13] Use Cases")
    slide_13_conclusion(prs)    ; print("  [13/13] Conclusion")

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "GuardianEye_FIM_Presentation.pptx")
    prs.save(out)
    print(f"\n  [OK]  Saved -> {out}")
    return out


if __name__ == "__main__":
    try:
        build_presentation()
    except ImportError:
        print("\n  [!] python-pptx not installed.")
        print("      Run:  pip install python-pptx")
        print("      Then: python make_ppt.py\n")
